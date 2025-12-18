"""
Text Extraction Module

This module contains functions for extracting text from various file formats.
It is used by Google Drive (MCP), WhatsApp, and file upload features to extract
text content from downloaded files.

Supported formats:
- PDF (via pdfplumber with table extraction)
- DOCX (via python-docx with table extraction)
- TXT (plain text)
- PPTX (via python-pptx)
- XLSX/XLS (via pandas)
- Images (via Tesseract OCR)

Two types of extractors are provided:
1. Stream-based extractors (for MCP client and in-memory processing)
2. File-path based extractors (for RAG system and uploaded files)
"""

import io
import re
import os
import PyPDF2
import pdfplumber
import pandas as pd
from docx import Document as DocxDocument
from pptx import Presentation
from PIL import Image
from pathlib import Path

# Check if Tesseract is available for OCR
try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False


# ==================== Helper Functions ====================

def deduplicate_repeated_chars(text):
    """Remove repeated characters (e.g., 'aaaa' -> 'a')"""
    if not text:
        return text
    return re.sub(r"(.)\1{2,}", r"\1", text)


def clean_pdf_text(text):
    """Clean and normalize PDF extracted text"""
    if not text:
        return text
    text = deduplicate_repeated_chars(text)
    text = re.sub(r" {2,}", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = "\n".join(line.strip() for line in text.split("\n"))
    return text.strip()


# ==================== Stream-based Extractors (for MCP client) ====================

def extract_text_from_pdf_stream(file_stream):
    """Extract text from PDF files using PyPDF2 (simple extraction)"""
    try:
        file_stream.seek(0)
        pdf_reader = PyPDF2.PdfReader(file_stream)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text.strip() if text.strip() else "No text content found in PDF"
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"


def extract_text_from_docx_stream(file_stream):
    """Extract text from DOCX files (simple extraction)"""
    try:
        file_stream.seek(0)
        doc = DocxDocument(file_stream)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return text.strip() if text.strip() else "No text content found in DOCX"
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"


def extract_text_from_txt_stream(file_stream):
    """Extract text from TXT files"""
    try:
        file_stream.seek(0)
        content = file_stream.read()
        if isinstance(content, bytes):
            return content.decode('utf-8', errors='ignore')
        return str(content)
    except Exception as e:
        return f"Error extracting TXT: {str(e)}"


def extract_text_from_pptx_stream(file_stream):
    """Extract text from PowerPoint files"""
    try:
        file_stream.seek(0)
        prs = Presentation(file_stream)
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n[Slide {i}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip() if text.strip() else "No text content found in PPTX"
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"


def extract_text_from_xlsx_stream(file_stream):
    """Extract text from Excel files"""
    try:
        file_stream.seek(0)
        df = pd.read_excel(file_stream, sheet_name=None)
        text = ""
        for sheet_name, sheet_data in df.items():
            text += f"\n[Sheet: {sheet_name}]\n"
            text += sheet_data.to_string(index=False) + "\n"
        return text.strip() if text.strip() else "No text content found in XLSX"
    except Exception as e:
        return f"Error extracting XLSX: {str(e)}"


def extract_text_from_image_stream(file_stream):
    """Extract text from images using OCR (Tesseract)"""
    if not TESSERACT_AVAILABLE:
        return "[OCR] Tesseract OCR not available. Please install pytesseract and tesseract-ocr."
    try:
        file_stream.seek(0)
        image = Image.open(file_stream)
        text = pytesseract.image_to_string(image)
        return f"[OCR Extracted]\n{text.strip()}" if text.strip() else "[OCR] No text found in image"
    except Exception as e:
        return f"Error extracting image text (OCR): {str(e)}"


def extract_text_by_mimetype(file_stream, mime_type, file_name=""):
    """Route to appropriate stream-based extractor based on MIME type"""
    extractors = {
        'application/pdf': extract_text_from_pdf_stream,
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': extract_text_from_docx_stream,
        'text/plain': extract_text_from_txt_stream,
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': extract_text_from_pptx_stream,
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': extract_text_from_xlsx_stream,
        'application/vnd.ms-excel': extract_text_from_xlsx_stream,
        'image/jpeg': extract_text_from_image_stream,
        'image/png': extract_text_from_image_stream,
        'image/jpg': extract_text_from_image_stream,
    }
    
    extractor = extractors.get(mime_type)
    if extractor:
        return extractor(file_stream)
    else:
        return f"Unsupported file type: {mime_type}"


# ==================== File-path based Extractors (for RAG system) ====================

def extract_pdf_with_tables(file_path):
    """
    Extract text from PDF files with table extraction using pdfplumber.
    This provides better table handling than PyPDF2.
    
    Args:
        file_path: Path to the PDF file
        
    Returns:
        Extracted text with tables formatted
    """
    parts = []

    try:
        with pdfplumber.open(file_path) as pdf:
            for page_number, page in enumerate(pdf.pages, start=1):
                page_parts = []
                text = page.extract_text()
                
                if text:
                    cleaned_text = clean_pdf_text(text)
                    if cleaned_text:
                        page_parts.append(cleaned_text)

                # Extract tables
                tables = page.extract_tables()
                for t_index, table in enumerate(tables):
                    if not table:
                        continue
                    rows_text = []
                    for row in table:
                        if row is None:
                            continue
                        cells = [
                            deduplicate_repeated_chars((cell or "").strip())
                            for cell in row
                        ]
                        rows_text.append("\t".join(cells))

                    if rows_text:
                        table_block = (
                            "\n----- TABLE -----\n"
                            + "\n".join(rows_text)
                            + "\n----- END TABLE -----\n"
                        )
                        page_parts.append(table_block)

                if page_parts:
                    parts.append("\n".join(page_parts))

        return "\n".join(parts)
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"


def extract_docx_with_tables(file_path):
    """
    Extract text from DOCX files with table extraction.
    Preserves document structure including tables.
    
    Args:
        file_path: Path to the DOCX file
        
    Returns:
        Extracted text with tables formatted
    """
    try:
        doc = DocxDocument(file_path)
        parts = []
        table_index = 0
        tables = doc.tables

        for element in doc.element.body:
            if element.tag.endswith("p"):
                paragraph = element.text.strip()
                if paragraph:
                    parts.append(paragraph)
            elif element.tag.endswith("tbl"):
                if table_index >= len(tables):
                    continue
                table = tables[table_index]
                table_index += 1

                rows_text = []
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells]
                    rows_text.append("\t".join(row_cells))

                table_block = (
                    "\n----- TABLE -----\n"
                    + "\n".join(rows_text)
                    + "\n----- END TABLE -----\n"
                )
                parts.append(table_block)

        return "\n".join(parts)
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"


def extract_txt_from_file(file_path):
    """Extract text from plain text files"""
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    except Exception as e:
        return f"Error extracting TXT: {str(e)}"


def extract_pptx_from_file(file_path):
    """Extract text from PowerPoint files"""
    try:
        prs = Presentation(file_path)
        text = ""
        for i, slide in enumerate(prs.slides, 1):
            text += f"\n[Slide {i}]\n"
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip() if text.strip() else "No text content found in PPTX"
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"


def extract_xlsx_from_file(file_path):
    """Extract text from Excel files"""
    try:
        df = pd.read_excel(file_path, sheet_name=None)
        text = ""
        for sheet_name, sheet_data in df.items():
            text += f"\n[Sheet: {sheet_name}]\n"
            text += sheet_data.to_string(index=False) + "\n"
        return text.strip() if text.strip() else "No text content found in XLSX"
    except Exception as e:
        return f"Error extracting XLSX: {str(e)}"


def extract_image_with_ocr(file_path):
    """Extract text from images using OCR (Tesseract)"""
    if not TESSERACT_AVAILABLE:
        return None, "Tesseract OCR not available"
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img, lang='eng')
        return text.strip() if text.strip() else None, None
    except Exception as e:
        return None, f"Error performing OCR: {str(e)}"


def extract_text_from_file(file_path):
    """
    Extract text from a file based on its extension.
    This is the main function for file-path based extraction.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Tuple of (extracted_text, file_type, error_message)
    """
    try:
        file_name = os.path.basename(file_path)
        suffix = Path(file_name).suffix.lower()
        
        if suffix == ".pdf":
            content = extract_pdf_with_tables(file_path)
            return content, "pdf", None
            
        elif suffix == ".docx":
            content = extract_docx_with_tables(file_path)
            return content, "docx", None
            
        elif suffix == ".txt":
            content = extract_txt_from_file(file_path)
            return content, "txt", None
            
        elif suffix == ".pptx":
            content = extract_pptx_from_file(file_path)
            return content, "pptx", None
            
        elif suffix in (".xlsx", ".xls"):
            content = extract_xlsx_from_file(file_path)
            return content, "excel", None
            
        elif suffix in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"):
            text, error = extract_image_with_ocr(file_path)
            if error:
                return None, "image", error
            return text, "ocr_image", None
            
        else:
            return None, None, f"Unsupported file type: {suffix}"
            
    except Exception as e:
        return None, None, f"Error extracting from {file_path}: {str(e)}"


# ==================== Utility Functions ====================

def is_tesseract_available():
    """Check if Tesseract OCR is available"""
    return TESSERACT_AVAILABLE

